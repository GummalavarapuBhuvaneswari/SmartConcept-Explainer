# ✅ AutoNote Transformer (Professional UI)

import streamlit as st
import fitz  # PyMuPDF
import google.generativeai as genai
from gtts import gTTS
from pptx import Presentation
import tempfile
import re
import os
import io
import time
from fpdf import FPDF
from pydub import AudioSegment
import base64


# === CONFIG ===
GEMINI_API_KEY = "AIzaSyAsoJ1zj2Q5YUnYBfkEG0bP6id6L-bqJPE"
genai.configure(api_key=GEMINI_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-flash")

# === CUSTOM CSS ===
def inject_custom_css():
    st.markdown("""
    <style>
        html, body, .main {
            background-color: var(--background-color);
            color: var(--text-color);
        }

        /* Headers */
        h1, h2, h3, h4, h5, h6 {
            color: var(--text-color);
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }

        /* Sidebar */
        .sidebar .sidebar-content {
            background-color: var(--background-color);
            box-shadow: 2px 0 10px rgba(0,0,0,0.1);
        }

        /* Buttons */
        .stButton>button {
            background-color: #3498db;
            color: white;
            border-radius: 5px;
            padding: 8px 16px;
            font-weight: 500;
            transition: all 0.3s;
            border: none;
        }
        .stButton>button:hover {
            background-color: #2980b9;
            transform: translateY(-1px);
            box-shadow: 0 2px 5px rgba(0,0,0,0.2);
        }

        /* Select boxes */
        .stSelectbox>div>div>select {
            border: 1px solid #dfe6e9;
            border-radius: 5px;
            padding: 8px;
        }

        /* File uploader */
        .stFileUploader>div>div {
            border: 2px dashed #bdc3c7;
            border-radius: 5px;
            padding: 20px;
        }

        /* Cards */
        .card {
            background: var(--secondary-background-color);
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            padding: 20px;
            margin-bottom: 20px;
            color: var(--text-color);
        }

        /* Divider */
        .divider {
            border-top: 1px solid #e0e0e0;
            margin: 20px 0;
        }

        /* Audio player description */
        .audio-caption {
            font-size: 14px;
            text-align: center;
            color: var(--text-color);
            margin-top: -10px;
        }
    </style>
    <script>
    const observer = new MutationObserver((mutations, obs) => {
        document.documentElement.style.setProperty('--background-color', getComputedStyle(document.body).backgroundColor);
        document.documentElement.style.setProperty('--secondary-background-color', getComputedStyle(document.body).backgroundColor === 'rgb(255, 255, 255)' ? '#ffffff' : '#1e1e1e');
        document.documentElement.style.setProperty('--text-color', getComputedStyle(document.body).color);
    });
    observer.observe(document.body, { attributes: true, childList: true, subtree: true });
    </script>
    """, unsafe_allow_html=True)

inject_custom_css()
# Initialize custom CSS
inject_custom_css()

# === PAGE SETUP ===
st.set_page_config(
    page_title="SmartConcept Explainer",
    page_icon="📘",
    layout="wide",
    initial_sidebar_state="expanded"
)

# === SESSION STATE ===
if 'concepts' not in st.session_state:
    st.session_state.concepts = []
if 'explanations' not in st.session_state:
    st.session_state.explanations = {}
if 'audio_files' not in st.session_state:
    st.session_state.audio_files = {}
if 'pdf_text' not in st.session_state:
    st.session_state.pdf_text = ""
if 'current_topic' not in st.session_state:
    st.session_state.current_topic = None
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

# === HEADER ===
with st.container():
    col1, col2 = st.columns([1, 3])
    with col1:
        st.image("https://cdn-icons-png.flaticon.com/512/2232/2232688.png", width=100)
    with col2:
        st.markdown("""
        <h1 style='margin-bottom: 0;'>SmartConcept Explainer</h1>
        <p style='color: #7f8c8d; margin-top: 0;'>AI-Powered Concept Explainer with Multilingual Support</p>
        """, unsafe_allow_html=True)
    
    st.markdown("---")

# === SIDEBAR ===
with st.sidebar:
    st.markdown("""
    <div style='text-align: center; margin-bottom: 30px;'>
        <h3 style='color: #2c3e50;'>📂 Document Upload</h3>
    </div>
    """, unsafe_allow_html=True)
    
    uploaded_file = st.file_uploader(
        "Upload your lecture notes (PDF/PPTX)",
        type=["pdf", "pptx"],
        label_visibility="collapsed"
    )
    
    st.markdown("---")
    
    st.markdown("""
    <div style='text-align: center;'>
        <h4 style='color: #2c3e50;'>⚙️ Settings</h4>
    </div>
    """, unsafe_allow_html=True)
    
    if 'concepts' in st.session_state and st.session_state.concepts:
        selected_lang = st.selectbox(
            "Explanation Language",
            ["English", "Telugu"],
            index=0,
            key="lang_select"
        )

# === TEXT EXTRACTORS ===
def extract_text_from_pdf(uploaded_file):
    with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
        return "\n".join(page.get_text() for page in doc)

def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# === CLEAN FOR AUDIO ===
def clean_for_voice(text, lang):
    text = re.sub(r"[#*>\\\-]", "", text)
    text = re.sub(r"\n+", ". ", text)
    text = re.sub(r"\s+", " ", text)

    if lang == "te":
        text = re.sub(r"(\d+)-(\d+)", r"\1 మైనస్ \2", text)
        text = re.sub(r"(\d+)-(\d+)", r"\1 టూ \2", text)
        text = re.sub(r"(\d+)\+(\d+)", r"\1 ప్లస్ \2", text)
        text = re.sub(r"(\d+)[*x×](\d+)", r"\1 మల్టీప్లైడ్ \2", text)
        text = re.sub(r"(\d+)[/÷](\d+)", r"\1 డివైడెడ్ బై \2", text)
        text = re.sub(r"(\d+)\^(\d+)", r"\1 పవర్ ఆఫ్ \2", text)
        text = re.sub(r"(\d+)/(\d+)", r"\1 ఆర్ \2", text)
        text = re.sub(r"(\d+)!(\d+)", r"\1 ఫ్యాక్టోరియల్ \2", text)
        text = text.replace("=", " సమానం ")
        text = text.replace(">", " ఎక్కువ")
        text = text.replace("<", " తక్కువ")
        text = re.sub(r"\d+", lambda m: " ".join(m.group()), text)
    return text.strip()

# === IDENTIFY CONCEPTS ===
def identify_concepts(text):
    """
    Smart concept extractor that identifies main topics and subtopics from any type of academic PDF.
    Works across math, science, and general texts.
    """
    prompt = f"""
You are a highly accurate document parser. Analyze the following academic content and extract only meaningful, high-level topics and subtopics.

STRICT INSTRUCTIONS:
- Only extract headings, chapter names, or subheadings — not formulas, not examples.
- Avoid garbage like short codes (e.g., TP, OM) or numeric-only lines.
- Do not return duplicates or repeated wording (e.g., "Permutation" and "Permutations").
- Do not include the chapter summary or index blocks (like Unit III...).
- Preserve the logical order from the text.

Return only a clean numbered list of the concepts found in the below text:

{text[:10000]}

FORMAT STRICTLY LIKE THIS:
1. Topic One
2. Topic Two
3. ...
"""

    try:
        response = gemini_model.generate_content(prompt)
        lines = response.text.strip().split("\n")

        clean = []
        for line in lines:
            if ". " in line:
                content = line.split(". ", 1)[1].strip()
                if len(content) > 5 and not re.fullmatch(r"[0-9. ]+", content):
                    clean.append(content)

        return clean

    except Exception as e:
        st.error(f"❌ Gemini extraction error: {e}")
        return []



# === EXPLAIN CONCEPT ===
def explain_concept(concept, context, lang):
    if lang == "Telugu":
        prompt = f"""
Explain the topic '{concept}' from the given PDF in a **simple, clear and student-friendly** manner using **Telugu + English mix** style.

Your explanation must follow this structure:

1. **Definition**  
   - Concept ni Telugu lo simple ga explain cheyyandi  
   - Important terminology English lo cheppandi (e.g., Mean అంటే Average)  
   - Example tho clarity ivvandi

2. **Key Characteristics**  
   - 5-6 important features cheppandi  
   - Each point simple sentence lo cheppandi  
   - Keywords English lo use cheyyandi

3. **Real-life Examples**  
   - PDF lo unna examples  
   - Mee side examples (3 total)  
   - Telugu explanation + English terms

4.  **Problems/Applications**  
   - Step-by-step ga problem solve cheyyandi  
   - English lo formula explain cheyyandi, Telugu lo steps explain cheyyandi  
   - Answer final ga cheppandi

5.  **Importance/Uses**  
   - Concept use enti, ela practical ga help chesthundi  
   - English terms explain cheyyandi (e.g., data analysis, decision making)

6.  **Common Mistakes**  
   - Students cheyyedhi common mistakes enti  
   - Examples tho explain cheyyandi

🧠 Style:
- Use **spoken conversational Telugu + English mix**
- Very simple language
- Highlight important English terms
- No lengthy or complex sentences
- Make it feel like a teacher is explaining to a student in class

Context:  
{context[:10000]}
        """
    else:
        prompt = f"""
Explain the topic '{concept}' from the given PDF in clear and detailed English.

Follow this structure:
1.  Definition
2.  Key Characteristics
3.  Real-life Examples
4.  Problems/Applications
5.  Importance/Uses
6.  Common Mistakes

- Use friendly and simple teaching style
- Each section should have at least 8–10 lines
- Include formulae, diagrams, and examples if found in PDF

Context:  
{context[:10000]}
        """

    try:
        response = gemini_model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"Explanation generation failed: {str(e)}")
        return f"Error generating explanation. Please try again. {str(e)}"


# === GENERATE AUDIO ===
# === AUDIO FUNCTIONS ===
def clean_telugu_text(text):
    """Enhanced Telugu text cleaning for better TTS output"""
    # Replace abbreviations and symbols
    replacements = {
        "ఉదా.": "ఉదాహరణకు",
        "అంటే": "అంటే ",
        "అనగా": "అనగా ",
        "-":"",
        "(": ", ",
        ")": ", ",
        "[": ", ",
        "]": ", ",
        "{": ", ",
        "}": ", ",
        "!":"",
        "•":"",
        "!":"ఫ్యాక్టోరియల్",
        "x":"ఇన్టూ",
        "=": " ఈక్వల్టూ ",
        "+": " ప్లస్ ",
        "-": " మైనస్ ",
        "/": " డివైడెడ్బై ",
        "∩":"ఇంటర్సెక్షన్",
        "∪":"యూనియన్",
        "÷":"డివైడెడ్బై",
        "Σ":"సమ్మషన్",
        "^": " పవర్ ",
        "%": " పర్సెంట్ ",
        "°C": " డిగ్రీల సెల్సియస్ ",
        "Dr.": "డాక్టర్ ",
        "Mr.": "మిస్టర్ ",
        "No.": "నంబర్ ",
        "e.g.": "ఉదాహరణకు ",
        "i.e.": "అంటే ",
        "Fig.": "ఫిగర్ ",
        "Eq.": "సమీకరణం ",
        "Ch.": "అధ్యాయం ",
        "*": "",        # Remove asterisks
        "#": "",        # Remove markdown headers
        "_": "",        # Remove underscores
        "\n": ". ",
        "\t": " ",
    }
    
    for k, v in replacements.items():
        text = text.replace(k, v)
    
    # Add pauses between sections
    text = re.sub(r"(\d+)\.", r"\n\n\1.", text)
    text = re.sub(r" +", " ", text)
    return text.strip()

def clean_english_text(text):
    """Enhanced English text cleaning for better TTS output"""
    replacements = {
        "e.g.": "for example",
        "i.e.": "that is",
        "Fig.": "Figure",
        "Eq.": "Equation",
        "Ch.": "Chapter",
        "*": "",        # Remove asterisks
        "#": "",        # Remove markdown headers
        "_": "",        # Remove underscores
        "\n": ". ",
        "\t": " ",
    }
    
    for k, v in replacements.items():
        text = text.replace(k, v)
    
    text = re.sub(r" +", " ", text)
    return text.strip()

def generate_high_quality_audio(text, lang):
    """Generate high quality audio for both Telugu and English"""
    try:
        # Configuration
        lang_code = "te" if lang == "Telugu" else "en"
        slow_speech = False  # Better clarity
        chunk_size = 1000 # Smaller chunks for better processing
        
        # Create temp directory
        temp_dir = tempfile.mkdtemp()
        final_path = os.path.join(temp_dir, f"audio_{hash(text)}.mp3")
        
        # Clean text based on language
        cleaned_text = clean_telugu_text(text) if lang == "Telugu" else clean_english_text(text)
        
        # Split into sentences first
        sentences = re.split(r'(?<=[.!?])\s+', cleaned_text)
        chunks = []
        current_chunk = ""
        
        # Create chunks of appropriate size
        for sentence in sentences:
            if len(current_chunk) + len(sentence) < chunk_size:
                current_chunk += " " + sentence
            else:
                chunks.append(current_chunk.strip())
                current_chunk = sentence
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Generate audio for each chunk
        audio_segments = []
        for i, chunk in enumerate(chunks):
            try:
                tts = gTTS(
                    text=chunk,
                    lang=lang_code,
                    slow=slow_speech,
                    lang_check=False  # Bypass strict language checking
                )
                chunk_path = os.path.join(temp_dir, f"chunk_{i}.mp3")
                tts.save(chunk_path)
                
                audio = AudioSegment.from_mp3(chunk_path)
                audio_segments.append(audio)
                
                # Add pause between chunks
                if i < len(chunks) - 1:
                    audio_segments.append(AudioSegment.silent(duration=300))
                
                os.remove(chunk_path)
            except Exception as e:
                st.warning(f"Could not process one audio segment: {str(e)}")
                continue
        
        # Combine all audio segments
        if audio_segments:
            combined_audio = audio_segments[0]
            for segment in audio_segments[1:]:
                combined_audio += segment
            
            # Export final audio
            combined_audio.export(final_path, format="mp3", bitrate="64k")
            return final_path
        
        return None
    
    except Exception as e:
        st.error(f"Audio generation failed: {str(e)}")
        return None


# === PDF DOWNLOAD ===
def download_explanations_as_pdf():
    import re
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Load and register Unicode font (required for Telugu)
    font_path = "NotoSans-Regular.ttf"
    if not os.path.exists(font_path):
        st.error("❌ Unicode font 'NotoSans-Regular.ttf' not found in the project directory.")
        return io.BytesIO()

    pdf.add_font("Noto", "", font_path, uni=True)
    pdf.set_font("Noto", size=12)

    for concept, explanation in st.session_state.explanations.items():
        # Set title in bold
        pdf.set_font("Noto", size=14)
        pdf.cell(0, 10, txt=concept, ln=True)
        pdf.set_font("Noto", size=12)

        # Clean up markdown-style text
        cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", explanation)  # remove bold markdown
        cleaned = re.sub(r"##+", "", cleaned)  # remove headings like ## Title
        cleaned = re.sub(r"`", "", cleaned)    # remove backticks
        cleaned = re.sub(r"\s+", " ", cleaned) # normalize spaces
        cleaned = re.sub(r"•", "-", cleaned)   # replace bullets if needed

        # Add text to PDF
        pdf.multi_cell(0, 8, cleaned + "\n")

        # Optional: Divider line
        pdf.ln(2)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(5)

    # Output PDF as bytes
    output = io.BytesIO()
    pdf.output(output)
    output.seek(0)
    return output



# === MAIN CONTENT ===
if uploaded_file and not st.session_state.pdf_text:
    with st.spinner("📄 Extracting content from document..."):
        ext = uploaded_file.name.split(".")[-1].lower()
        if ext == "pdf":
            st.session_state.pdf_text = extract_text_from_pdf(uploaded_file)
        elif ext == "pptx":
            st.session_state.pdf_text = extract_text_from_pptx(uploaded_file)
    
    with st.spinner("🧠 Analyzing document for key concepts..."):
        st.session_state.concepts = identify_concepts(st.session_state.pdf_text)
        if st.session_state.concepts:
            st.success("✅ Document processed successfully!")
            with st.expander("📋 Extracted Concepts", expanded=True):
                cols = st.columns(2)
                for i, concept in enumerate(st.session_state.concepts):
                    with cols[i % 2]:
                        st.info(f"• {concept}")

# === CONCEPT EXPLANATION SECTION ===
if st.session_state.concepts:
    st.markdown("---")

    with st.container():
        st.markdown("""
        <div style='text-align: center; margin-bottom: 20px;'>
            <h2>Concept Explanation</h2>
        </div>
        """, unsafe_allow_html=True)

        selected_topic = st.selectbox(
            "Select a concept to explain:",
            st.session_state.concepts,
            key="topic_select"
        )

        col1, col2 = st.columns(2)
        with col1:
            if st.button("Generate Explanation", key="explain_btn", use_container_width=True):
                with st.spinner(f"Generating {selected_lang} explanation..."):
                    explanation = explain_concept(
                        selected_topic,
                        st.session_state.pdf_text,
                        selected_lang  # ✅ Use selected_lang here instead of undefined 'lang'
                    )
                    st.session_state.explanations[selected_topic] = explanation

        with col2:
            if st.button("Generate Audio", key="audio_btn", use_container_width=True):
                if selected_topic in st.session_state.explanations:
                    with st.spinner("Generating audio explanation..."):
                        audio_path = generate_high_quality_audio(
                            st.session_state.explanations[selected_topic], selected_lang
                        )
                        st.session_state.audio_files[selected_topic] = audio_path
                else:
                    st.warning("Please generate the explanation first")

        if selected_topic in st.session_state.explanations:
            st.markdown("---")
            with st.container():
                st.markdown(f"""
                <div class='card'>
                    <h3 style='color: #3498db;'>{selected_topic}</h3>
                    <div class='divider'></div>
                    {st.session_state.explanations[selected_topic]}
                </div>
                """, unsafe_allow_html=True)

                if selected_topic in st.session_state.audio_files and st.session_state.audio_files[selected_topic]:
                    st.markdown(f"### 🔊 {selected_lang} Audio: {selected_topic}")
                    audio_file_path = st.session_state.audio_files[selected_topic]
                    audio_bytes = open(audio_file_path, 'rb').read()
                    st.audio(audio_bytes, format='audio/mp3')
                    # Add caption below player
                    st.markdown("""
                    <div style="font-size: 14px; text-align: center; color: var(--text-color); margin-top: -10px;">
                        🎧 <em>Playback speed: 1.25x recommended for better clarity</em>
                    </div>
                    """, unsafe_allow_html=True)


# === EXPORT SECTION ===
if st.session_state.explanations:
    st.markdown("---")
    with st.container():
        st.markdown("""
        <div style='text-align: center; margin-bottom: 20px;'>
            <h2>Export Options</h2>
        </div>
        """, unsafe_allow_html=True)
        
        pdf_bytes = download_explanations_as_pdf()
        st.download_button(
            label="Download All Explanations as PDF",
            data=pdf_bytes,
            file_name="AutoNote_Explanations.pdf",
            mime="application/pdf",
            use_container_width=True
        )

# === CHAT SECTION ===
# === CHAT SECTION ===
st.markdown("---")
with st.container():
    st.markdown("""
    <div style='text-align: center; margin-bottom: 20px;'>
        <h2>🧠 Concept Assistant</h2>
        <p>Ask questions about any concept (Chat Memory Enabled)</p>
    </div>
    """, unsafe_allow_html=True)

    # Display chat history
    for chat in st.session_state.chat_history:
        with st.chat_message("user", avatar="🧑‍💻"):
            st.markdown(chat["user"])
        with st.chat_message("assistant", avatar="🤖"):
            st.markdown(chat["assistant"])

    # Chat input box
    user_question = st.chat_input("Type your question here...")
    if user_question:
        with st.chat_message("user", avatar="🧑‍💻"):
            st.markdown(user_question)

        with st.chat_message("assistant", avatar="🤖"):
            chat_model = genai.GenerativeModel("gemini-1.5-flash")
            response = chat_model.generate_content(f"Context: {st.session_state.pdf_text[:2000]}\n\nQuestion: {user_question}")
            st.markdown(response.text)

            # Save to chat history
            st.session_state.chat_history.append({
                "user": user_question,
                "assistant": response.text
            })
